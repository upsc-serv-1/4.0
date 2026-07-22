import os
import pptx
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

def create_animated_pptx(input_pptx_path, output_pptx_path):
    prs = pptx.Presentation(input_pptx_path)

    for slide_idx, slide in enumerate(prs.slides):
        # Gather content shape IDs (cards, text boxes below header title)
        animated_shapes = []
        for shape in slide.shapes:
            if shape.top and shape.top.inches > 1.2 and shape.top.inches < 6.8:
                animated_shapes.append(shape)

        if not animated_shapes:
            continue

        # Build clean OpenXML p:timing block
        # PowerPoint standard sequential click entrance animation XML
        seq_items = []
        for i, shape in enumerate(animated_shapes):
            node_id = (i + 1) * 2
            item_xml = f"""<p:cTn id="{node_id}" fill="hold" nodeType="clickEffect">
              <p:stCondLst>
                <p:cond delay="0"/>
              </p:stCondLst>
              <p:childTnLst>
                <p:set>
                  <p:cb>
                    <p:spTgt spid="{shape.shape_id}"/>
                  </p:cb>
                  <p:to>
                    <p:animVal>
                      <p:strVal val="visible"/>
                    </p:animVal>
                  </p:to>
                </p:set>
              </p:childTnLst>
            </p:cTn>"""
            seq_items.append(item_xml)

        combined_items = "".join(seq_items)
        first_spid = animated_shapes[0].shape_id
        last_spid = animated_shapes[-1].shape_id

        timing_xml = f"""<p:timing {nsdecls('p')}>
          <p:tnLst>
            <p:par id="1" nodeType="root">
              <p:childTnLst>
                <p:seq id="2" concurrent="1" nextAc="seek">
                  <p:cTn id="3" dur="indefinite" nodeType="mainSeq">
                    <p:childTnLst>
                      {combined_items}
                    </p:childTnLst>
                  </p:cTn>
                  <p:prevCondLst>
                    <p:cond evt="onPrev" delay="0">
                      <p:tgtEl>
                        <p:spTgt spid="{first_spid}"/>
                      </p:tgtEl>
                    </p:cond>
                  </p:prevCondLst>
                  <p:nextCondLst>
                    <p:cond evt="onNext" delay="0">
                      <p:tgtEl>
                        <p:spTgt spid="{last_spid}"/>
                      </p:tgtEl>
                    </p:cond>
                  </p:nextCondLst>
                </p:seq>
              </p:childTnLst>
            </p:par>
          </p:tnLst>
        </p:timing>"""

        # Remove existing timing node if any
        existing_timing = slide.element.find(pptx.oxml.ns.qn('p:timing'))
        if existing_timing is not None:
            slide.element.remove(existing_timing)

        timing_elem = parse_xml(timing_xml)
        slide.element.append(timing_elem)

    prs.save(output_pptx_path)
    print(f"Animated PPTX generated successfully: {output_pptx_path}")

if __name__ == "__main__":
    input_path = r"C:\Users\Dr. Yogesh\Downloads\Sociology_Optional_Orientation_Updated.pptx"
    output_path = r"c:\Users\Dr. Yogesh\Videos\APP FOLDER - V1 - Copy\app\frontend-noji-2.6.2\3\pilot pro 10.2\Sociology_Optional_Orientation_Animated.pptx"
    create_animated_pptx(input_path, output_path)
