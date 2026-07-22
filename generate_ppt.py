import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_sociology_presentation_light(output_path):
    prs = Presentation()
    # 16:9 Widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Elegant Light Theme Color Palette
    COLOR_BG = RGBColor(248, 250, 252)        # Light Slate / Soft Off-White (#F8FAFC)
    COLOR_CARD = RGBColor(255, 255, 255)      # Pure White Card (#FFFFFF)
    COLOR_CARD_BORDER = RGBColor(226, 232, 240) # Subtle Border (#E2E8F0)
    COLOR_PRIMARY = RGBColor(37, 99, 235)     # Royal Blue (#2563EB)
    COLOR_ACCENT = RGBColor(217, 119, 6)      # Warm Amber Gold (#D97706)
    COLOR_TEXT_MAIN = RGBColor(15, 23, 42)    # Deep Slate Text (#0F172A)
    COLOR_TEXT_MUTED = RGBColor(71, 85, 105)  # Soft Slate Gray (#475569)
    COLOR_TEAL = RGBColor(13, 148, 136)       # Teal Accent (#0D9488)

    def add_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="UPSC CSE SOCIOLOGY OPTIONAL"):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_ACCENT
        p_cat.font.name = "Calibri"

        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.75))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_MAIN
        p_title.font.name = "Georgia"

    def add_card(slide, left, top, width, height, title, items):
        # Card Shape
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = COLOR_CARD_BORDER
        card.line.width = Pt(1.5)

        # Card Title
        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.55))
        tf = tb.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(17)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_PRIMARY
        p0.font.name = "Georgia"

        # Items
        tb_content = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.65), width - Inches(0.4), height - Inches(0.75))
        tf_c = tb_content.text_frame
        tf_c.word_wrap = True

        for idx, item in enumerate(items):
            p = tf_c.add_paragraph() if idx > 0 else tf_c.paragraphs[0]
            p.space_after = Pt(10)
            p.font.size = Pt(13)
            p.font.name = "Calibri"

            if isinstance(item, tuple):
                run1 = p.add_run()
                run1.text = "• " + item[0] + ": "
                run1.font.bold = True
                run1.font.color.rgb = COLOR_TEAL

                run2 = p.add_run()
                run2.text = item[1]
                run2.font.color.rgb = COLOR_TEXT_MUTED
            else:
                run = p.add_run()
                run.text = "• " + item
                run.font.color.rgb = COLOR_TEXT_MUTED

    blank_layout = prs.slide_layouts[6]

    # ==================== SLIDE 1: TITLE SLIDE ====================
    slide1 = prs.slides.add_slide(blank_layout)
    add_background(slide1)
    
    main_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.0), Inches(11.333), Inches(5.5))
    main_card.fill.solid()
    main_card.fill.fore_color.rgb = COLOR_CARD
    main_card.line.color.rgb = COLOR_PRIMARY
    main_card.line.width = Pt(2)

    tb = slide1.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "UPSC CSE MAINS • SOCIOLOGY OPTIONAL"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    p.space_after = Pt(12)

    p2 = tf.add_paragraph()
    p2.text = "Sociology Optional: Strategy & Roadmap"
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT_MAIN
    p2.font.name = "Georgia"
    p2.space_after = Pt(8)

    p3 = tf.add_paragraph()
    p3.text = "300+ Marks Strategy, Concept Clarity & Mentorship Plan"
    p3.font.size = Pt(18)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_PRIMARY
    p3.space_after = Pt(20)

    p4 = tf.add_paragraph()
    p4.text = "Key Pillars: Concept Clarity | Answer Writing | Thinker Linkages | PYQ Structuring"
    p4.font.size = Pt(13)
    p4.font.color.rgb = COLOR_TEXT_MUTED
    p4.space_after = Pt(20)

    p5 = tf.add_paragraph()
    p5.text = "Mentor: UPSC Mains Candidate & Sociology Mentor  |  Batch Orientation Session"
    p5.font.size = Pt(12)
    p5.font.bold = True
    p5.font.color.rgb = COLOR_TEXT_MAIN

    slide1.notes_slide.notes_text_frame.text = "Welcome students. Explain 300+ score vision and demystify Sociology Optional."

    # ==================== SLIDE 2: WHY SOCIOLOGY? ====================
    slide2 = prs.slides.add_slide(blank_layout)
    add_background(slide2)
    add_header(slide2, "Why Sociology Optional? (Kyun Chune Sociology?)")

    col_w = Inches(3.6)
    gap = Inches(0.4)
    left_m = Inches(0.8)
    top_m = Inches(1.7)
    h = Inches(5.0)

    add_card(slide2, left_m, top_m, col_w, h, "1. High Scoring (300+ Target)", [
        ("Predictable PYQs", "Repeated question themes make preparation laser-targeted."),
        ("70% Static Syllabus", "Finish core theoretical concepts before Prelims."),
        ("Objective Evaluation", "Thinker-anchored answers get top marks consistently.")
    ])

    add_card(slide2, left_m + col_w + gap, top_m, col_w, h, "2. Massive GS & Essay Overlap", [
        ("GS Paper 1 (Society)", "~75 marks directly covered (Caste, Women, Urbanization)."),
        ("125-Mark Essay", "Provides strong sociological frameworks for social essays."),
        ("GS Paper 4 (Ethics)", "Enriches understanding of values, norms & case studies.")
    ])

    add_card(slide2, left_m + (col_w + gap)*2, top_m, col_w, h, "3. Easy & Accessible", [
        ("No Prior Background", "Engineering, Science & Commerce students excel equally."),
        ("Relatable Real Reality", "Study family, caste, gender, and social movements."),
        ("Compact Footprint", "Smaller syllabus footprint compared to History/Geography.")
    ])

    slide2.notes_slide.notes_text_frame.text = "Explain ROI of Sociology in Mains. Address myth: Non-humanities students can easily score 300+."

    # ==================== SLIDE 3: WHAT IS SOCIOLOGY ====================
    slide3 = prs.slides.add_slide(blank_layout)
    add_background(slide3)
    add_header(slide3, "Syllabus Breakdown (Sociology Me Kya Padhte Hain?)")

    w2 = Inches(5.6)
    gap2 = Inches(0.5)

    add_card(slide3, left_m, top_m, w2, h, "Paper 1: Fundamentals of Sociology", [
        ("Sociology as Science", "Methodology, Positivism vs Interpretive paradigms."),
        ("6 Core Western Thinkers", "Marx, Durkheim, Weber, Parsons, Merton, Mead."),
        ("Social Stratification", "Concepts of Class, Status, Gender & Social Mobility."),
        ("Social Institutions", "Work & Economy, Politics, Religion, Family & Kinship.")
    ])

    add_card(slide3, left_m + w2 + gap2, top_m, w2, h, "Paper 2: Indian Society Structure & Change", [
        ("Perspectives on India", "Indology (Ghurye), Structural-Functionalism (Srinivas), Marxist (Desai)."),
        ("Social Structure", "Caste system, agrarian class, untouchability, tribal issues."),
        ("Social Transformation", "Modernization, Green Revolution, Peasant/Tribal movements."),
        ("Paper 1 + 2 Linkage", "Applying Paper 1 thinkers' theories to Indian social realities.")
    ])

    slide3.notes_slide.notes_text_frame.text = "Define C. Wright Mills' Sociological Imagination. Highlight Paper 1 theory + Paper 2 application."

    # ==================== SLIDE 4: ANSWER WRITING MASTERY ====================
    slide4 = prs.slides.add_slide(blank_layout)
    add_background(slide4)
    add_header(slide4, "Pillar #1: Answer Writing Mastery (Answer Writing Kaise Karein?)")

    add_card(slide4, left_m, top_m, col_w, h, "1. Anatomy of 250+ Answer", [
        ("Introduction (10-15%)", "Sociological definition or thinker reference."),
        ("Body (75-80%)", "Multi-perspective arguments, thinker quotes, empirical data."),
        ("Conclusion (10-15%)", "Synthesis of views & futuristic policy alignment.")
    ])

    add_card(slide4, left_m + col_w + gap, top_m, col_w, h, "2. Managing 10 / 15 / 20 Markers", [
        ("10-Marker (150 words)", "Punchy, direct, 2 core thinkers + clean concept definition."),
        ("15-Marker (250 words)", "Multidimensional points, Indian context + counter-views."),
        ("20-Marker (250-300w)", "Deep critical analysis, thinker comparisons & synthesis.")
    ])

    add_card(slide4, left_m + (col_w + gap)*2, top_m, col_w, h, "3. Value Additions", [
        ("Key Terminology", "Anomie, Hegemony, Sanskritization, Little/Great Traditions."),
        ("Thinker Anchoring", "Avoid plain GS style; anchor arguments with sociologists."),
        ("Empirical Data", "Quote NFHS, Census, NITI Aayog & EPW studies.")
    ])

    slide4.notes_slide.notes_text_frame.text = "Explain difference between GS answer and Sociology answer. Focus on terminology and thinker quotes."

    # ==================== SLIDE 5: BRAINSTORMING ====================
    slide5 = prs.slides.add_slide(blank_layout)
    add_background(slide5)
    add_header(slide5, "Pillar #2: 360° Brainstorming (60 Seconds Mind Mapping)")

    add_card(slide5, left_m, top_m, col_w, h, "1. 3 Sociological Lenses", [
        ("Functionalist Lens", "Social order, consensus, harmony, and integration."),
        ("Conflict / Marxist Lens", "Power dynamics, class struggle, exploitation & inequality."),
        ("Feminist Lens", "Gender impact, patriarchy, unpaid labor, intersectionality.")
    ])

    add_card(slide5, left_m + col_w + gap, top_m, col_w, h, "2. Analytical Frameworks", [
        ("PESTLE Model", "Political, Economic, Social, Tech, Legal, Environmental angles."),
        ("Micro vs. Macro", "Individual agency (Mead/Goffman) vs. Macro structures (Marx/Parsons)."),
        ("Tradition vs Modernity", "Continuity of traditional norms alongside modern institutions.")
    ])

    add_card(slide5, left_m + (col_w + gap)*2, top_m, col_w, h, "3. 60-Second Mind Map Drill", [
        ("Deconstruct Keyword", "Understand 'Critically Examine', 'Elucidate', 'Discuss'."),
        ("Draft 5 Points", "Write 5 quick keywords in margin before writing full answer."),
        ("Visual Presentation", "Use mini-tables, side headings, and underlined key terms.")
    ])

    slide5.notes_slide.notes_text_frame.text = "Teach 60-second mind map formula to avoid getting stuck during exam."

    # ==================== SLIDE 6: INTERLINKAGES ====================
    slide6 = prs.slides.add_slide(blank_layout)
    add_background(slide6)
    add_header(slide6, "Pillar #3: Content Structuring & Interlinkages (Paper 1 + 2 Linkage)")

    add_card(slide6, left_m, top_m, w2, h, "Paper 1 + Paper 2 Interlinkage Matrix", [
        ("Karl Marx ➔ Indian Agrarian Distress", "Alienation & class struggle applied to Green Revolution inequality."),
        ("Max Weber ➔ Indian Administration", "Rational-legal authority vs patrimonialism in governance."),
        ("Durkheim ➔ Anomie & Farmer Suicides", "Breakdown of traditional social solidarity in modern times."),
        ("Feminist Thinkers ➔ Dalit Feminism", "Sharmila Rege's Dalit standpoint theory in Indian caste matrix.")
    ])

    add_card(slide6, left_m + w2 + gap2, top_m, w2, h, "Dynamic Content & Visual Presentation", [
        ("Current Affairs Integration", "Link daily news (Caste Census, Gig Workers, Urbanization) to theory."),
        ("Visual Flowcharts & Tables", "Use thinker comparison tables, Venn diagrams, and process cycles."),
        ("Empirical Field Studies", "Cite Srinivas (Rampura), Beteille (Sripuram), Gail Omvedt, Sujata Patel."),
        ("Clean Paragraphing", "Structured 3-4 line paragraphs with bold key concepts.")
    ])

    slide6.notes_slide.notes_text_frame.text = "Highlight that 300+ marks toppers stand out because of seamless Paper 1 + Paper 2 interlinkage."

    # ==================== SLIDE 7: SMART NOTE-MAKING ====================
    slide7 = prs.slides.add_slide(blank_layout)
    add_background(slide7)
    add_header(slide7, "Pillar #4: Smart Note-Making Strategy (Smart Notes Kaise Banayein?)")

    add_card(slide7, left_m, top_m, col_w, h, "1. One-Page Thinker Sheet", [
        ("Core Methodology", "Positivist / Interpretive / Critical stance in 2 lines."),
        ("Key Theories & Terms", "Top 4-5 concepts with exact definitions."),
        ("Major Critiques", "Counter-thinkers and flaws."),
        ("Mains PYQ Utility", "Where to use this thinker in PYQs.")
    ])

    add_card(slide7, left_m + col_w + gap, top_m, col_w, h, "2. Value Addition Register", [
        ("Data & Fact Book", "NFHS-5, Census, NITI Aayog & NCRB statistics."),
        ("Sociologists & Books", "Index of famous books and authors."),
        ("Government Schemes", "Sociological evaluation of welfare policies.")
    ])

    add_card(slide7, left_m + (col_w + gap)*2, top_m, col_w, h, "3. PYQ-Centric Notes", [
        ("10-Year PYQ Mapping", "Group past questions topic-wise."),
        ("Ready Intro-Conclusions", "Templates for recurring sub-themes."),
        ("50-Page Micro-Notes", "Ultra-condensed rapid revision notes before Mains.")
    ])

    slide7.notes_slide.notes_text_frame.text = "Warn against copying textbooks. Explain 1-Page Thinker Sheet & 50-Page micro-notes strategy."

    # ==================== SLIDE 8: BATCH METHODOLOGY ====================
    slide8 = prs.slides.add_slide(blank_layout)
    add_background(slide8)
    add_header(slide8, "Batch Learning Methodology (Hum Kaise Padhenge?)")

    add_card(slide8, left_m, top_m, col_w, h, "1. Zero-to-Hero Lectures", [
        ("Conceptual Clarity", "Start from NCERT foundation up to Haralambos & Ritzer level."),
        ("Bilingual (Hinglish)", "Easy Hinglish explanation with real-life Indian examples."),
        ("Thinker Masterclasses", "Dedicated deep-dive sessions on all core thinkers.")
    ])

    add_card(slide8, left_m + col_w + gap, top_m, col_w, h, "2. Daily & Weekly Drills", [
        ("Daily 1-Question Practice", "Build daily writing habit without burnout."),
        ("Live Answer Deconstruction", "Break down model answers live in class."),
        ("Brainstorming Workshops", "Weekly live drills on complex PYQs.")
    ])

    add_card(slide8, left_m + (col_w + gap)*2, top_m, col_w, h, "3. Curated Deliverables", [
        ("Crisp Mind Maps", "Handouts eliminating heavy multi-book reading."),
        ("10-Year PYQ Bank", "Solved model answers for 10 years' questions."),
        ("Current Sociology Digest", "Monthly compilation of sociological current affairs.")
    ])

    slide8.notes_slide.notes_text_frame.text = "Outline classroom workflow: Concept Lecture -> Live PYQ Discussion -> Student Writing -> Personal Evaluation."

    # ==================== SLIDE 9: ROADMAP ====================
    slide9 = prs.slides.add_slide(blank_layout)
    add_background(slide9)
    add_header(slide9, "Mentorship Roadmap (3-Phase Action Plan)")

    add_card(slide9, left_m, top_m, col_w, h, "Phase 1: Foundation (W1-8)", [
        ("Focus", "Paper 1 Fundamentals & 6 Western Thinkers."),
        ("Activity", "Concept building, note making & daily 1 answer."),
        ("Milestone", "Complete mastery over theoretical foundation.")
    ])

    add_card(slide9, left_m + col_w + gap, top_m, col_w, h, "Phase 2: Application (W9-16)", [
        ("Focus", "Paper 2 Indian Society & Paper 1-2 interlinkages."),
        ("Activity", "Sectional Tests + Value addition data integration."),
        ("Milestone", "100% syllabus coverage & PYQ mapping.")
    ])

    add_card(slide9, left_m + (col_w + gap)*2, top_m, col_w, h, "Phase 3: Mains Drill", [
        ("Focus", "Full Length Tests under real exam conditions."),
        ("Activity", "1-on-1 Mentorship evaluation & speed optimization."),
        ("Milestone", "300+ Marks Readiness in Mains!")
    ])

    slide9.notes_slide.notes_text_frame.text = "Walk through 3-phase timeline. Highlight 1-on-1 personal mentorship."

    # ==================== SLIDE 10: ACTION PLAN ====================
    slide10 = prs.slides.add_slide(blank_layout)
    add_background(slide10)
    add_header(slide10, "Day 1 Action Plan (Aapka Pehla Kadam)")

    add_card(slide10, left_m, top_m, col_w, h, "1. Immediate Day 1 Tasks", [
        ("Print Syllabus", "Keep printed syllabus on study table; read keywords daily."),
        ("NCERT Reading", "Read Class 11th & 12th Sociology NCERTs carefully."),
        ("PYQ Scanning", "Browse past 5 years' papers to understand question style.")
    ])

    add_card(slide10, left_m + col_w + gap, top_m, col_w, h, "2. Recommended Booklist", [
        ("Paper 1 Base", "Haralambos & Holborn (Essential/Orange) or Ritzer (Thinkers)."),
        ("Paper 2 Base", "Social Change in Modern India (MN Srinivas) & Ram Ahuja."),
        ("Class Handouts", "Our curated consolidated handouts will be primary!")
    ])

    add_card(slide10, left_m + (col_w + gap)*2, top_m, col_w, h, "3. Golden Mantra & Q&A", [
        ("Consistency Over Intensity", "Write 1 answer daily rather than 10 once a month."),
        ("Trust the Process", "Sociology is rewarding when approached logically."),
        ("Open Floor for Doubts", "Let's address all your questions and doubts now!")
    ])

    slide10.notes_slide.notes_text_frame.text = "Give Day 1 homework. Open floor for student Q&A."

    # Save presentation
    prs.save(output_path)
    print(f"Light Theme Presentation saved to: {output_path}")

if __name__ == "__main__":
    out_dir = r"C:\Users\Dr. Yogesh\.gemini\antigravity\brain\109a927a-a0d0-416e-90b8-71167717f5b0"
    os.makedirs(out_dir, exist_ok=True)
    ppt_path = os.path.join(out_dir, "Sociology_Optional_Introductory_Class.pptx")
    create_sociology_presentation_light(ppt_path)
