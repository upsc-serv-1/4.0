import pptx
from pptx.enum.dml import MSO_FILL

prs = pptx.Presentation(r'C:\Users\Dr. Yogesh\Downloads\Sociology_Optional_Orientation_Updated.pptx')
print(f"Slide Dimensions: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f} inches (16:9)")

for idx, slide in enumerate(prs.slides):
    print(f"\n==================== SLIDE {idx+1} ====================")
    for shape_idx, shape in enumerate(slide.shapes):
        name = shape.name
        left = shape.left.inches
        top = shape.top.inches
        width = shape.width.inches
        height = shape.height.inches
        
        color_str = ""
        try:
            if shape.fill.type == MSO_FILL.SOLID:
                color_str = f"SOLID RGB:{shape.fill.fore_color.rgb}"
        except Exception:
            pass

        print(f"Shape #{shape_idx+1}: '{name}' | Pos: (L={left:.2f}\", T={top:.2f}\", W={width:.2f}\", H={height:.2f}\") | {color_str}")
        
        if shape.has_text_frame:
            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                txt = p.text.strip()
                if txt:
                    font_name = p.font.name if (p.font and p.font.name) else "Inherited"
                    font_size = p.font.size.pt if (p.font and p.font.size) else "Inherited"
                    c_rgb = ""
                    try:
                        if p.font and p.font.color and p.font.color.rgb:
                            c_rgb = str(p.font.color.rgb)
                    except Exception:
                        pass
                    print(f"   Para: '{txt[:50]}' [Font: {font_name}, Size: {font_size}, Color: {c_rgb}]")
